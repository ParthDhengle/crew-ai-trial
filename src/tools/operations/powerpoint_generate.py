# powerpoint_generate.py

import os
import sys
import tempfile
from datetime import datetime
import warnings
import json
import re
from dotenv import load_dotenv
import platform
import subprocess

# Suppress syntax warnings
warnings.filterwarnings("ignore", category=SyntaxWarning)

# Absolute imports to avoid relative import issues
try:
    from utils.logger import setup_logger
    from common_functions.Find_project_root import find_project_root
except ImportError:
    def setup_logger():
        import logging
        logger = logging.getLogger("AIAssistant")
        logger.setLevel(logging.DEBUG)
        if not logger.handlers:
            console_handler = logging.StreamHandler()
            console_handler.setLevel(logging.INFO)
            formatter = logging.Formatter(
                "%(asctime)s - %(name)s - %(levelname)s - %(message)s",
                datefmt="%Y-%m-%d %H:%M:%S"
            )
            console_handler.setFormatter(formatter)
            logger.addHandler(console_handler)
        return logger
    
    def find_project_root():
        import os
        current_dir = os.path.abspath(os.path.dirname(__file__))
        root_dir = current_dir
        while root_dir != os.path.dirname(root_dir):
            if os.path.exists(os.path.join(root_dir, '.env')):  # Use '.git', 'requirements.txt', or another root marker if preferred
                return root_dir
            root_dir = os.path.dirname(root_dir)
        return current_dir  # Fallback to script's directory if no marker found

logger = setup_logger()
PROJECT_ROOT = find_project_root()

# Load .env file dynamically from project root
env_path = os.path.join(PROJECT_ROOT, '.env')
if not os.path.exists(env_path):
    logger.error(f".env file not found at {env_path}")
    sys.exit(1)
load_dotenv(env_path)
logger.info(f"Loaded .env file from {env_path}")

def check_powerpoint_installation():
    """Check if Microsoft PowerPoint is installed on the system."""
    try:
        if platform.system() == "Windows":
            # Common PowerPoint installation paths
            possible_paths = [
                r"C:\Program Files\Microsoft Office\root\Office16\POWERPNT.EXE",
                r"C:\Program Files (x86)\Microsoft Office\root\Office16\POWERPNT.EXE",
            ]
            
            for path in possible_paths:
                if os.path.exists(path):
                    logger.info(f"Microsoft PowerPoint found at: {path}")
                    return True, path
            
            logger.warning("Microsoft PowerPoint not found in common installation paths")
            return False, None
        else:
            logger.warning("Microsoft PowerPoint auto-open is only supported on Windows")
            return False, None
    except Exception as e:
        logger.error(f"Error checking PowerPoint installation: {str(e)}")
        return False, None

def call_grok(prompt: str, api_key: str) -> str:
    """Call Grok API directly."""
    try:
        from groq import Groq
        client = Groq(api_key=api_key)
        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}],
            max_tokens=1500,
            temperature=0.3
        )
        return response.choices[0].message.content
    except Exception as e:
        logger.warning(f"Grok API failed: {str(e)}")
        return None

def call_gemini(prompt: str, api_key: str) -> str:
    """Call Gemini API as fallback."""
    try:
        import google.generativeai as genai
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel("gemini-1.5-flash-latest")
        response = model.generate_content(prompt)
        return response.text
    except Exception as e:
        logger.warning(f"Gemini API failed: {str(e)}")
        return None

def clean_llm_response(response: str) -> str:
    """Clean LLM response to extract JSON."""
    if not response:
        return None
    
    logger.debug(f"Raw LLM response: {response[:500]}...")
    
    # Remove markdown code blocks
    response = re.sub(r'^```json\s*|\s*```$', '', response.strip(), flags=re.MULTILINE)
    response = re.sub(r'^```\s*|\s*```$', '', response.strip(), flags=re.MULTILINE)
    response = response.strip()
    
    # Try to find JSON within the response
    json_match = re.search(r'\{.*\}', response, re.DOTALL)
    if json_match:
        json_str = json_match.group(0)
        logger.debug(f"Extracted JSON: {json_str[:200]}...")
        return json_str
    
    logger.warning(f"No JSON found in response: {response[:200]}...")
    return None

def create_fallback_presentation_content(query: str) -> dict:
    """Create fallback presentation content when LLM fails."""
    logger.info("Creating fallback presentation content")
    
    return {
        "title": f"Presentation for: {query}",
        "slides": [
            {
                "heading": "Title Slide",
                "content": ["This is a basic presentation generated as fallback.", "Please provide a more detailed query for better results."]
            },
            {
                "heading": "Summary",
                "content": ["No specific content generated due to LLM failure."]
            }
        ]
    }

def create_powerpoint_presentation(output_path: str, content: dict):
    """Create a Microsoft PowerPoint presentation using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        pres = Presentation()
        
        # Add title slide
        slide_layout = pres.slide_layouts[0]  # Title slide layout
        slide = pres.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = content.get("title", "Generated Presentation")
        
        # Add content slides
        for section in content.get("slides", []):
            slide_layout = pres.slide_layouts[1]  # Title and Content layout
            slide = pres.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = section.get("heading", "Slide")
            
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            for point in section.get("content", []):
                p = tf.add_paragraph()
                p.text = point
        
        pres.save(output_path)
        logger.info(f"PowerPoint presentation created at: {output_path}")
        return True
    except ImportError:
        logger.error("python-pptx library not found. Please install it with 'pip install python-pptx'")
        return False
    except Exception as e:
        logger.error(f"Error creating PowerPoint presentation: {str(e)}")
        return False

def powerpoint_generate_from_query(query: str) -> tuple[bool, str]:
    """
    Generates a Microsoft PowerPoint presentation based on a user query using AI-driven content generation.
    
    Args:
        query (str): User query describing desired presentation (e.g., "Create a presentation on climate change").
    
    Returns:
        tuple[bool, str]: (success, result message or error)
    """
    try:
        # Check PowerPoint installation
        ppt_installed, ppt_path = check_powerpoint_installation()
        if not ppt_installed:
            logger.warning("Microsoft PowerPoint not found. Presentation will be created but may not open automatically.")
        
        # Verify API keys
        grok_api_key = os.getenv("GROQ_API_KEY3")
        gemini_api_key = os.getenv("GEMINI_API_KEY3")
        if not grok_api_key and not gemini_api_key:
            logger.error(f"No API keys found for GROQ_API_KEY1 or GEMINI_API_KEY1 in {env_path}")
            return False, f"Error: No API keys found for GROQ_API_KEY1 or GEMINI_API_KEY1 in {env_path}"

        logger.info(f"Generating PowerPoint presentation for query: {query}")
        
        # Construct improved prompt for LLM
        prompt = f"""
You are a presentation expert. Generate content for a PowerPoint presentation based on this request.

User Query: "{query}"

Return ONLY a valid JSON object with this exact structure:
{{
  "title": "Presentation Title",
  "slides": [
    {{
      "heading": "Slide Heading",
      "content": ["Bullet point 1", "Bullet point 2"]
    }}
  ]
}}

Make the content informative and relevant to the query.
Generate 5-10 slides with concise bullet points.
Return only the JSON, no explanation.
"""
        
        # Call LLM with retry logic
        response = None
        content = None
        
        for attempt in range(3):
            logger.info(f"Attempting LLM call #{attempt + 1}")
            
            if grok_api_key and not response:
                response = call_grok(prompt, grok_api_key)
                if response:
                    cleaned = clean_llm_response(response)
                    if cleaned:
                        try:
                            content = json.loads(cleaned)
                            logger.info("Successfully parsed Grok response")
                            break
                        except json.JSONDecodeError as e:
                            logger.warning(f"Grok JSON decode error: {str(e)}")
                            response = None
            
            if gemini_api_key and not content:
                response = call_gemini(prompt, gemini_api_key)
                if response:
                    cleaned = clean_llm_response(response)
                    if cleaned:
                        try:
                            content = json.loads(cleaned)
                            logger.info("Successfully parsed Gemini response")
                            break
                        except json.JSONDecodeError as e:
                            logger.warning(f"Gemini JSON decode error: {str(e)}")
                            response = None
        
        # Use fallback if LLM failed
        if not content:
            logger.warning("LLM parsing failed after all retries. Using fallback content.")
            content = create_fallback_presentation_content(query)
        
        logger.info(f"Presentation content: {content}")
        
        # Create temporary file for presentation
        output_dir = tempfile.mkdtemp(prefix="powerpoint_presentation_")
        presentation_name = f"auto_presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        presentation_path = os.path.join(output_dir, presentation_name)
        
        if not create_powerpoint_presentation(presentation_path, content):
            return False, "Failed to create PowerPoint presentation. Ensure python-pptx is installed."
        
        success_message = f"PowerPoint presentation created at: {presentation_path}"
        
        if ppt_installed:
            try:
                # Try to open with PowerPoint if available
                os.startfile(presentation_path)
                success_message += "\nPresentation opened in Microsoft PowerPoint."
            except Exception as e:
                logger.warning(f"Failed to open presentation: {str(e)}")
                success_message += f"\nCould not auto-open presentation: {str(e)}"
        else:
            # Open the directory instead if PowerPoint not found
            if platform.system() == "Windows":
                os.startfile(output_dir)
            success_message += "\nPresentation directory opened."
        
        logger.info(success_message)
        return True, success_message
        
    except Exception as e:
        logger.error(f"Error generating PowerPoint presentation: {str(e)}")
        return False, f"Error: {str(e)}"

if __name__ == "__main__":
    # Ensure environment variables are loaded
    grok_api_key = os.getenv("GROQ_API_KEY3")
    gemini_api_key = os.getenv("GEMINI_API_KEY3")
    if not grok_api_key and not gemini_api_key:
        print(f"Error: Please set GROQ_API_KEY1 or GEMINI_API_KEY1 in {env_path}")
        sys.exit(1)
    
    # Example usage
    test_query = "Create a detailed presentation on the impacts of artificial intelligence on society."
    success, result = powerpoint_generate_from_query(test_query)
    print(f"Success: {success}")
    print(result)
